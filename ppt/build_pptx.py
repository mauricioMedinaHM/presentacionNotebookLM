#!/usr/bin/env python3
"""
Genera "Amplificador_Cognitivo.pptx" replicando la presentación React
(app/src/components/slides/Slide00..18) con tema oscuro, colores, tipografías
y ANIMACIONES de entrada equivalentes a las del código (GSAP).

Las animaciones GSAP se mapean a animaciones de entrada nativas de PowerPoint,
inyectando el árbol <p:timing> en el XML de cada slide. Reproducen el orden y
los retardos (delays) originales y se disparan automáticamente al entrar a la
diapositiva (igual que en la web).

Mapa de animaciones:
  FadeUp / BlurReveal / ScaleIn / plain  -> Fade
  SlideIn(right)                         -> Fly In desde la derecha
  ClipReveal                             -> Wipe hacia arriba
  WipeIn                                 -> Wipe hacia la derecha (revela izq->der)
  GlitchReveal                           -> Fly In desde la izquierda
  StaggerReveal                          -> Fade escalonado por item
"""

import copy
from pathlib import Path
from lxml import etree
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap

# ---------------------------------------------------------------------------
# Rutas / assets
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
PUBLIC = ROOT / "app" / "public"
OUT = Path(__file__).resolve().parent / "Amplificador_Cognitivo.pptx"
TMP = Path(__file__).resolve().parent / ".assets"
TMP.mkdir(exist_ok=True)

# ---------------------------------------------------------------------------
# Paleta (de index.css)
# ---------------------------------------------------------------------------
BASE    = "07070D"
SURFACE = "111118"
WHITE   = "FFFFFF"
ACCENT  = "FF6A2A"   # lava
AMBER   = "FFB800"
PRIMARY = "FF3B3B"   # rojo
SUBTLE  = "C0C0D0"
GREY    = "CCCCCC"
GREY2   = "BBBBBB"
MUTED   = "8A8A9A"

DISPLAY = "Urbanist"
SANS    = "Plus Jakarta Sans"
MONO    = "JetBrains Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Imágenes (reescala la foto pesada)
# ---------------------------------------------------------------------------
def prep_photo(src, dst, max_w=900):
    im = Image.open(src).convert("RGB")
    if im.width > max_w:
        h = int(im.height * max_w / im.width)
        im = im.resize((max_w, h), Image.LANCZOS)
    im.save(dst, "JPEG", quality=86)
    return dst

PHOTO = prep_photo(PUBLIC / "fotoMauricio.jpg", TMP / "foto.jpg")
QR = PUBLIC / "qrIgMauri.png"

# ---------------------------------------------------------------------------
# Presentación base
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def pct_x(p): return int(SLIDE_W * p)
def pct_y(p): return int(SLIDE_H * p)
def px(n):    return Inches(n / 96.0)   # CSS px -> inch


# ---------------------------------------------------------------------------
# Builder de una diapositiva
# ---------------------------------------------------------------------------
class SlideBuilder:
    def __init__(self):
        self.slide = prs.slides.add_slide(BLANK)
        self.anims = []          # {spid, effect, delay, dur}
        self._bg()

    # -- fondo oscuro --
    def _bg(self):
        bg = self.slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(BASE)

    # -- orbe de brillo (ellipse translúcida) --
    def glow(self, color, size_px, x, y, alpha=10):
        d = px(size_px)
        sp = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(SLIDE_W * x - d / 2), int(SLIDE_H * y - d / 2), d, d,
        )
        sp.line.fill.background()
        sp.shadow.inherit = False
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(color)
        # alpha
        srgb = sp.fill.fore_color._xFill.find(qn("a:srgbClr"))
        a = etree.SubElement(srgb, qn("a:alpha"))
        a.set("val", str(int(alpha * 1000)))
        # mandar al fondo
        spTree = sp._element.getparent()
        spTree.remove(sp._element)
        spTree.insert(2, sp._element)
        return sp

    # -- caja de texto --
    def text(self, x, y, w, h, paras, anchor="t", anim=None, delay=0.0,
             dur=700, wrap=True):
        tb = self.slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                              "b": MSO_ANCHOR.BOTTOM}[anchor]
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = para.get("align", PP_ALIGN.LEFT)
            if "space_after" in para:
                p.space_after = Pt(para["space_after"])
            if "line" in para:
                p.line_spacing = para["line"]
            for run in para["runs"]:
                r = p.add_run()
                r.text = run["t"]
                r.font.size = Pt(run.get("sz", 18))
                r.font.bold = run.get("bold", False)
                r.font.italic = run.get("italic", False)
                r.font.name = run.get("font", SANS)
                r.font.color.rgb = RGBColor.from_string(run.get("color", SUBTLE))
                spc = run.get("spc")
                if spc is not None:
                    r.font._rPr.set("spc", str(int(spc * 100)))
        if anim:
            self.register(tb, anim, delay, dur)
        return tb

    def image(self, path, x, y, w=None, h=None, anim=None, delay=0.0, dur=700,
              round_corners=False):
        pic = self.slide.shapes.add_picture(str(path), x, y, w, h)
        if round_corners:
            pic._element.spPr.append(
                etree.fromstring(
                    '<a:prstGeom xmlns:a="%s" prst="roundRect"><a:avLst>'
                    '<a:gd name="adj" fmla="val 8000"/></a:avLst></a:prstGeom>' % A
                )
            )
        if anim:
            self.register(pic, anim, delay, dur)
        return pic

    def chip(self, cx, top, label, color, anim=None, delay=0.0):
        """etiqueta tipo 'CONCEPTO CLAVE' centrada en cx (eje x relativo)."""
        return self.text(
            int(SLIDE_W * cx - Inches(3)), top, Inches(6), Inches(0.5),
            [{"align": PP_ALIGN.CENTER, "runs": [
                {"t": label, "sz": 14, "bold": True, "color": color,
                 "font": SANS, "spc": 2.5}]}],
            anim=anim, delay=delay,
        )

    def bar(self, x, y, w, h, color, alpha=100, radius=True):
        sp = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            x, y, w, h)
        sp.line.fill.background()
        sp.shadow.inherit = False
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(color)
        if alpha < 100:
            srgb = sp.fill.fore_color._xFill.find(qn("a:srgbClr"))
            etree.SubElement(srgb, qn("a:alpha")).set("val", str(int(alpha * 1000)))
        return sp

    def register(self, shape, effect, delay, dur):
        self.anims.append({"spid": shape.shape_id, "effect": effect,
                           "delay": int(delay * 1000), "dur": int(dur)})

    # -- inyecta el <p:timing> con todas las animaciones --
    def finalize(self):
        if not self.anims:
            return
        EFFECTS = {
            "fade":   ("fade",            10, 0),
            "flyR":   ("slide(fromRight)", 2, 2),
            "flyL":   ("slide(fromLeft)",  2, 8),
            "flyB":   ("slide(fromBottom)",2, 4),
            "wipeU":  ("wipe(up)",        22, 4),
            "wipeR":  ("wipe(right)",     22, 2),
        }
        cid = [10]
        def nid():
            cid[0] += 1
            return cid[0]

        effects_xml = []
        for a in self.anims:
            filt, pid, sub = EFFECTS[a["effect"]]
            e_main = nid(); e_set = nid(); e_eff = nid()
            effects_xml.append(f"""
            <p:par>
              <p:cTn id="{e_main}" presetID="{pid}" presetClass="entr" presetSubtype="{sub}" fill="hold" grpId="0" nodeType="withEffect">
                <p:stCondLst><p:cond delay="{a['delay']}"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="{e_set}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                      <p:tgtEl><p:spTgt spid="{a['spid']}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:animEffect transition="in" filter="{filt}">
                    <p:cBhvr>
                      <p:cTn id="{e_eff}" dur="{a['dur']}"/>
                      <p:tgtEl><p:spTgt spid="{a['spid']}"/></p:tgtEl>
                    </p:cBhvr>
                  </p:animEffect>
                </p:childTnLst>
              </p:cTn>
            </p:par>""")

        timing = f"""<p:timing xmlns:p="{P}" xmlns:a="{A}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>{''.join(effects_xml)}
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""
        self.slide._element.append(etree.fromstring(timing))


# ---------------------------------------------------------------------------
# helpers de runs
# ---------------------------------------------------------------------------
def run(t, sz, color=SUBTLE, bold=False, italic=False, font=SANS, spc=None):
    return {"t": t, "sz": sz, "color": color, "bold": bold, "italic": italic,
            "font": font, "spc": spc}

def para(runs, align=PP_ALIGN.CENTER, space_after=None, line=None):
    d = {"runs": runs, "align": align}
    if space_after is not None: d["space_after"] = space_after
    if line is not None: d["line"] = line
    return d

CENTER = PP_ALIGN.CENTER
LEFT = PP_ALIGN.LEFT
RIGHT = PP_ALIGN.RIGHT


# ===========================================================================
# SLIDES
# ===========================================================================

# ---- Slide 00 : portada -----------------------------------------------------
def slide00():
    s = SlideBuilder()
    s.glow(ACCENT, 500, 0.15, 0.50, 9)
    s.glow(AMBER, 400, 0.80, 0.30, 8)
    # dots
    s.text(Inches(1.0), Inches(1.4), Inches(3), Inches(0.4),
           [para([run("● ● ●", 12, ACCENT)], LEFT)], anim="fade", delay=0.2)
    # título
    s.text(Inches(1.0), Inches(1.9), Inches(7.0), Inches(2.6),
           [para([run("El ", 54, WHITE, bold=True, font=DISPLAY),
                  run("Amplificador", 54, ACCENT, bold=True, font=DISPLAY)], LEFT),
            para([run("Cognitivo", 54, WHITE, bold=True, font=DISPLAY)], LEFT)],
           anim="wipeU", delay=0.4, dur=1000)
    # subtítulo
    s.text(Inches(1.0), Inches(4.35), Inches(6.5), Inches(1.0),
           [para([run("Cómo usar IA para estudiar, investigar y documentar mejor.",
                      20, SUBTLE)], LEFT, line=1.25)],
           anim="fade", delay=0.8)
    # nombre
    s.text(Inches(1.0), Inches(5.5), Inches(6.5), Inches(0.5),
           [para([run("—  ", 16, ACCENT),
                  run("MAURICIO MEDINA", 16, ACCENT, bold=True, spc=2.0)], LEFT)],
           anim="fade", delay=1.1)
    # foto
    s.image(PHOTO, Inches(9.0), Inches(1.6), w=Inches(3.2),
            anim="flyR", delay=0.5, dur=900, round_corners=True)
    # hint
    s.text(Inches(0), Inches(6.85), SLIDE_W, Inches(0.4),
           [para([run("USA LAS FLECHAS PARA NAVEGAR  ⌄", 11, "B0B0C0", spc=2.0)], CENTER)],
           anim="fade", delay=1.5)
    s.finalize()


# ---- Slide 01 : ¿quién estudia? --------------------------------------------
def slide01():
    s = SlideBuilder()
    s.glow(ACCENT, 600, 0.50, 0.50, 9)
    s.text(Inches(0), Inches(1.0), SLIDE_W, Inches(0.8),
           [para([run("✋", 40, ACCENT)], CENTER)], anim="fade", delay=0.2)
    s.text(Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.0),
           [para([run("¿Quién está estudiando algo ", 46, WHITE, bold=True, font=DISPLAY),
                  run("en este momento", 46, ACCENT, bold=True, font=DISPLAY),
                  run("?", 46, WHITE, bold=True, font=DISPLAY)], CENTER, line=1.05)],
           anim="wipeU", delay=0.4, dur=1000)
    s.text(Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.8),
           [para([run("Carrera, curso, o por su propia cuenta", 26, SUBTLE)], CENTER)],
           anim="fade", delay=0.8)
    s.text(Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.6),
           [para([run("Carrera", 20, AMBER), run("   /   ", 20, MUTED),
                  run("Curso", 20, AMBER), run("   /   ", 20, MUTED),
                  run("Autodidacta", 20, AMBER)], CENTER)],
           anim="fade", delay=1.2)
    s.finalize()


# ---- Slide 02 : contexto ----------------------------------------------------
def slide02():
    s = SlideBuilder()
    s.glow(AMBER, 500, 0.25, 0.40, 8)
    s.glow(ACCENT, 400, 0.75, 0.60, 8)
    s.chip(0.5, Inches(1.1), "◆  CONCEPTO CLAVE", ACCENT, anim="fade", delay=0.2)
    s.text(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.6),
           [para([run("La importancia del ", 46, WHITE, bold=True, font=DISPLAY),
                  run("contexto", 46, ACCENT, bold=True, font=DISPLAY)], CENTER, line=1.05)],
           anim="wipeU", delay=0.4, dur=1000)
    s.text(Inches(2.0), Inches(4.0), Inches(9.3), Inches(1.1),
           [para([run("Acá todos vamos a usar IA, pero busquemos entender lo que "
                      "estamos haciendo con ella.", 22, SUBTLE)], CENTER, line=1.25)],
           anim="fade", delay=0.7)
    s.text(Inches(2.0), Inches(5.4), Inches(9.3), Inches(1.2),
           [para([run("La IA va a trabajar bajo el contexto que le podamos dar.",
                      26, AMBER, bold=True)], CENTER, line=1.2)],
           anim="fade", delay=1.0, dur=1000)
    s.finalize()


# ---- Slide 03 : segundo cerebro --------------------------------------------
def slide03():
    s = SlideBuilder()
    s.glow(AMBER, 500, 0.50, 0.40, 8)
    s.text(Inches(0.9), Inches(2.0), Inches(6.6), Inches(2.0),
           [para([run("NotebookLM como nuestro ", 42, WHITE, bold=True, font=DISPLAY),
                  run("segundo cerebro", 42, ACCENT, bold=True, font=DISPLAY)], LEFT, line=1.08)],
           anim="wipeU", delay=0.4, dur=1000)
    s.text(Inches(0.9), Inches(4.4), Inches(6.4), Inches(1.6),
           [para([run("En esta IA podemos meter cientos de libros, videos e imágenes "
                      "y lo va a guardar en su contexto.", 22, SUBTLE)], LEFT, line=1.3)],
           anim="fade", delay=0.7)
    # "órbita" estática: cerebro central + nodos
    cx, cy = Inches(10.3), Inches(3.75)
    s.text(int(cx - Inches(1)), int(cy - Inches(0.5)), Inches(2), Inches(1),
           [para([run("🧠", 48, ACCENT)], CENTER)], anim="fade", delay=0.5)
    nodes = [("📖", -1.7, 0), ("🎬", 0, -1.5), ("🖼️", 1.7, 0), ("📄", 0, 1.5)]
    for i, (ic, dx, dy) in enumerate(nodes):
        s.text(int(cx + Inches(dx) - Inches(0.5)), int(cy + Inches(dy) - Inches(0.4)),
               Inches(1), Inches(0.8),
               [para([run(ic, 24, WHITE)], CENTER)], anim="fade", delay=0.6 + i * 0.12)
    s.finalize()


# ---- Slide 04 : ALUCINACIÓN ------------------------------------------------
def slide04():
    s = SlideBuilder()
    s.glow(PRIMARY, 600, 0.50, 0.50, 10)
    # orbe icono
    d = px(112)
    orb = s.slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   int(SLIDE_W / 2 - d / 2), Inches(0.7), d, d)
    orb.line.color.rgb = RGBColor.from_string(PRIMARY)
    orb.line.width = Pt(1.5)
    orb.shadow.inherit = False
    orb.fill.solid(); orb.fill.fore_color.rgb = RGBColor.from_string(PRIMARY)
    srgb = orb.fill.fore_color._xFill.find(qn("a:srgbClr"))
    etree.SubElement(srgb, qn("a:alpha")).set("val", "8000")
    s.register(orb, "fade", 0.2, 700)
    s.text(int(SLIDE_W / 2 - d / 2), Inches(0.85), d, px(80),
           [para([run("⚡", 30, PRIMARY)], CENTER, line=1.0)], anchor="m",
           anim="fade", delay=0.2)
    s.chip(0.5, Inches(2.0), "CONCEPTO CLAVE", PRIMARY, anim="fade", delay=0.4)
    s.text(Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.4),
           [para([run("ALUCINACIÓN", 68, WHITE, bold=True, font=DISPLAY, spc=0.5)], CENTER)],
           anim="flyL", delay=0.6, dur=700)
    s.text(Inches(2.0), Inches(4.35), Inches(9.3), Inches(1.1),
           [para([run("Cuando la IA inventa información y te la dice con tanta "
                      "seguridad que parece verdad.", 22, SUBTLE)], CENTER, line=1.25)],
           anim="fade", delay=0.9)
    s.text(Inches(2.0), Inches(5.7), Inches(9.3), Inches(1.0),
           [para([run("La IA no sabe que no sabe. Por eso miente con buena gramática.",
                      20, GREY, italic=True)], CENTER, line=1.25)],
           anim="fade", delay=1.2, dur=1000)
    s.finalize()


# ---- Slide 05 : ¿cuánto miente? (barras) -----------------------------------
def slide05():
    s = SlideBuilder()
    s.glow(ACCENT, 400, 0.15, 0.30, 8)
    s.glow(AMBER, 400, 0.85, 0.70, 8)
    s.text(Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.9),
           [para([run("¿Cuánto miente la IA que usás todos los días?",
                      34, WHITE, bold=True, font=DISPLAY)], CENTER)],
           anim="wipeU", delay=0.2, dur=900)
    s.text(Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.5),
           [para([run("Tasa de alucinación en alta complejidad", 16, GREY2)], CENTER)],
           anim="fade", delay=0.35)
    models = [
        ("Claude (Anthropic)", "3.7 Sonnet / Opus", "15% – 22%", 18.5, ACCENT),
        ("ChatGPT (OpenAI)", "GPT-4.1 / 4o", "18% – 25%", 21.5, AMBER),
        ("Gemini (Google)", "2.0 Pro", "20% – 28%", 24.0, PRIMARY),
        ("DeepSeek", "V3 / R1", "24% – 32%", 28.0, PRIMARY),
    ]
    top = Inches(2.15)
    rowh = Inches(1.05)
    for i, (name, model, rng, pctv, color) in enumerate(models):
        y = int(top + rowh * i)
        # fila bg
        rowbg = s.bar(Inches(0.8), y, Inches(11.7), Inches(0.9), SURFACE, alpha=35)
        # textos
        s.text(Inches(1.1), y + Inches(0.12), Inches(3.2), Inches(0.7),
               [para([run(name, 16, WHITE, bold=True)], LEFT, space_after=2),
                para([run(model, 13, GREY2)], LEFT)],
               anim="fade", delay=0.4 + i * 0.15)
        # track
        track = s.bar(Inches(4.5), y + Inches(0.32), Inches(5.6), Inches(0.26), "0A0A12")
        # fill
        fw = Inches(5.6) * (pctv * 2.5 / 100.0)
        s.bar(Inches(4.5), y + Inches(0.32), int(fw), Inches(0.26), color)
        s.text(Inches(10.4), y + Inches(0.22), Inches(2.0), Inches(0.5),
               [para([run(rng, 16, GREY, font=MONO)], RIGHT)],
               anim="fade", delay=0.45 + i * 0.15)
    s.text(Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.6),
           [para([run("Ninguna está libre. La pregunta no es si se equivoca, es cuánto.",
                      16, GREY2)], CENTER)], anim="fade", delay=1.2)
    s.finalize()


# ---- Slide 06 : 22% – 94% ---------------------------------------------------
def slide06():
    s = SlideBuilder()
    s.glow(PRIMARY, 500, 0.50, 0.40, 9)
    s.chip(0.5, Inches(0.8), "△  EL PROBLEMA ES REAL", PRIMARY, anim="fade", delay=0.2)
    s.text(Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.6),
           [para([run("22%", 90, WHITE, bold=True, font=DISPLAY),
                  run("   —   ", 44, "5A5A6A"),
                  run("94%", 90, ACCENT, bold=True, font=DISPLAY)], CENTER)],
           anim="fade", delay=0.5, dur=800)
    s.text(Inches(2.0), Inches(3.35), Inches(9.3), Inches(1.0),
           [para([run("Tasa de alucinación en los 26 modelos más usados del mundo "
                      "en tareas exigentes.", 20, SUBTLE)], CENTER, line=1.25)],
           anim="fade", delay=0.9)
    s.text(Inches(2.0), Inches(4.45), Inches(9.3), Inches(0.5),
           [para([run("Stanford AI Index 2026", 15, ACCENT, font=MONO, spc=1.5)], CENTER)],
           anim="fade", delay=1.0)
    # caja inferior
    box = s.bar(Inches(3.0), Inches(5.2), Inches(7.3), Inches(1.55), SURFACE, alpha=35)
    box.line.color.rgb = RGBColor.from_string(PRIMARY); box.line.width = Pt(1.25)
    s.text(Inches(3.2), Inches(5.45), Inches(6.9), Inches(0.7),
           [para([run("+362", 24, WHITE, bold=True),
                  run("  incidentes documentados de IA en 2025", 18, GREY)], CENTER)],
           anim="fade", delay=1.2, dur=1000)
    s.text(Inches(3.2), Inches(6.15), Inches(6.9), Inches(0.5),
           [para([run("55% más que en 2024", 17, AMBER)], CENTER)],
           anim="fade", delay=1.3)
    s.finalize()


# ---- Slide 07 : CONTEXTO ----------------------------------------------------
def slide07():
    s = SlideBuilder()
    s.glow(AMBER, 600, 0.50, 0.50, 9)
    s.chip(0.5, Inches(1.0), "▦  CONCEPTO CLAVE", AMBER, anim="fade", delay=0.2)
    s.text(Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.4),
           [para([run("CONTEXTO", 68, WHITE, bold=True, font=DISPLAY, spc=1.0)], CENTER)],
           anim="flyL", delay=0.4, dur=700)
    s.text(Inches(2.0), Inches(3.35), Inches(9.3), Inches(1.1),
           [para([run("Todo lo que la IA tiene a mano para responderte en ese momento.",
                      26, SUBTLE)], CENTER, line=1.25)],
           anim="fade", delay=0.7)
    s.text(Inches(2.0), Inches(5.0), Inches(3.8), Inches(1.2),
           [para([run("Sin contexto", 20, PRIMARY, bold=True)], CENTER, space_after=6),
            para([run("la IA adivina", 17, GREY)], CENTER)],
           anim="fade", delay=1.0)
    s.text(Inches(6.0), Inches(5.25), Inches(1.3), Inches(0.6),
           [para([run("→", 28, "999999")], CENTER)], anim="fade", delay=1.0)
    s.text(Inches(7.5), Inches(5.0), Inches(3.8), Inches(1.2),
           [para([run("Con contexto", 20, AMBER, bold=True)], CENTER, space_after=6),
            para([run("la IA razona", 17, GREY)], CENTER)],
           anim="fade", delay=1.1)
    s.finalize()


# ---- Slide 08 : biblioteca o 3 libros --------------------------------------
def slide08():
    s = SlideBuilder()
    s.glow(ACCENT, 400, 0.25, 0.60, 8)
    s.glow(AMBER, 400, 0.75, 0.40, 8)
    s.text(Inches(2.0), Inches(1.0), Inches(9.3), Inches(1.4),
           [para([run("Si tu profe te da 3 libros con todas las respuestas...",
                      32, SUBTLE)], CENTER, line=1.2)],
           anim="fade", delay=0.2, dur=1000)
    s.text(Inches(1.5), Inches(3.2), Inches(3.8), Inches(2.2),
           [para([run("📚", 56, ACCENT)], CENTER, space_after=14),
            para([run("¿Vas a la biblioteca entera?", 22, ACCENT, bold=True)],
                 CENTER, line=1.15)],
           anim="fade", delay=0.5)
    s.text(Inches(6.0), Inches(3.9), Inches(1.3), Inches(0.8),
           [para([run("o", 28, GREY)], CENTER)], anim="fade", delay=0.7)
    s.text(Inches(8.0), Inches(3.2), Inches(3.8), Inches(2.2),
           [para([run("🔖", 56, AMBER)], CENTER, space_after=14),
            para([run("¿O te quedás con esos 3?", 22, AMBER, bold=True)],
                 CENTER, line=1.15)],
           anim="fade", delay=0.9)
    s.finalize()


# ---- Slide 09 : 67% reducción ----------------------------------------------
def slide09():
    s = SlideBuilder()
    s.glow(ACCENT, 500, 0.50, 0.40, 9)
    s.text(Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.7),
           [para([run("Las alucinaciones no se eliminan por completo", 20, SUBTLE)], CENTER)],
           anim="fade", delay=0.2)
    s.text(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.6),
           [para([run("↘  ", 40, ACCENT),
                  run("67%", 80, ACCENT, bold=True, font=DISPLAY)], CENTER)],
           anim="fade", delay=0.5, dur=800)
    s.text(Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.6),
           [para([run("de reducción", 24, WHITE, bold=True)], CENTER)],
           anim="fade", delay=0.6)
    # barra
    s.bar(Inches(4.0), Inches(4.6), Inches(5.3), Inches(0.16), SURFACE)
    s.bar(Inches(4.0), Inches(4.6), int(Inches(5.3) * 0.33), Inches(0.16), ACCENT)
    s.text(Inches(4.0), Inches(4.85), Inches(5.3), Inches(0.5),
           [para([run("Sin NotebookLM", 14, GREY)], LEFT)], anim="fade", delay=0.9)
    s.text(Inches(4.0), Inches(4.85), Inches(5.3), Inches(0.5),
           [para([run("Con NotebookLM", 14, GREY)], RIGHT)], anim="fade", delay=0.9)
    s.text(Inches(2.5), Inches(5.7), Inches(8.3), Inches(1.2),
           [para([run("Pueden pensar que no se le tiene que dar importancia, pero el "
                      "día de mañana ese porcentaje de error puede costar muy caro...",
                      18, GREY)], CENTER, line=1.3)],
           anim="fade", delay=1.3, dur=1000)
    s.finalize()


# ---- Slide 10 : abrir NotebookLM + notepad ---------------------------------
def slide10():
    s = SlideBuilder()
    s.glow(AMBER, 500, 0.50, 0.50, 8)
    s.text(Inches(0), Inches(0.9), SLIDE_W, Inches(0.8),
           [para([run("🖥️", 36, AMBER)], CENTER)], anim="fade", delay=0.2)
    s.text(Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.6),
           [para([run("Abramos ", 42, WHITE, bold=True, font=DISPLAY),
                  run("NotebookLM", 42, ACCENT, bold=True, font=DISPLAY),
                  run(" y nuestro notepad preferido", 42, WHITE, bold=True, font=DISPLAY)],
                 CENTER, line=1.08)],
           anim="wipeU", delay=0.4, dur=1000)
    s.text(Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.7),
           [para([run("Puede ser Notion, Google Docs o Obsidian", 22, SUBTLE)], CENTER)],
           anim="fade", delay=0.7)
    tools = [("✎  Notion", 0.0), ("⌘  Google Docs", 0.12), ("✦  Obsidian", 0.24)]
    xs = [Inches(3.0), Inches(5.7), Inches(8.4)]
    for (label, dl), x in zip(tools, xs):
        s.text(x, Inches(4.7), Inches(2.6), Inches(0.6),
               [para([run(label, 20, WHITE, bold=True)], CENTER)],
               anim="fade", delay=1.0 + dl)
    s.text(Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
           [para([run("Lo más importante es que interprete archivos Markdown",
                      17, GREY2)], CENTER)], anim="fade", delay=1.3)
    s.finalize()


# ---- Slide 11 : PREGUNTAR TODO ---------------------------------------------
def slide11():
    s = SlideBuilder()
    s.glow(ACCENT, 700, 0.50, 0.50, 10)
    s.text(Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.0),
           [para([run("PREGUNTAR", 64, WHITE, bold=True, font=DISPLAY)], CENTER, line=1.05),
            para([run("ABSOLUTAMENTE", 64, ACCENT, bold=True, font=DISPLAY)], CENTER, line=1.05),
            para([run("TODO.", 64, WHITE, bold=True, font=DISPLAY)], CENTER, line=1.05)],
           anim="flyL", delay=0.3, dur=700)
    # divider
    s.bar(int(SLIDE_W / 2 - Inches(1.5)), Inches(4.9), Inches(3.0), Pt(2), ACCENT)
    s.text(Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.0),
           [para([run("NO DEJAR UNA SOLA DUDA.", 36, SUBTLE, bold=True, font=DISPLAY)],
                 CENTER)],
           anim="fade", delay=1.3, dur=800)
    s.finalize()


# ---- Slide 12 : EL FLUJO ---------------------------------------------------
def slide12():
    s = SlideBuilder()
    s.glow(ACCENT, 400, 0.10, 0.30, 8)
    s.glow(AMBER, 400, 0.90, 0.70, 8)
    s.text(Inches(1.0), Inches(0.7), Inches(11.3), Inches(1.1),
           [para([run("EL ", 56, WHITE, bold=True, font=DISPLAY),
                  run("FLUJO", 56, ACCENT, bold=True, font=DISPLAY)], CENTER)],
           anim="wipeU", delay=0.2, dur=900)
    steps = [
        ("1.", "⬆", "Cargo las fuentes en NotebookLM"),
        ("2.", "?", "Pregunto TODO lo que no entiendo"),
        ("3.", "✎", "Anoto con mis palabras en Obsidian"),
        ("4.", "🔗", "Enlazo conceptos entre sí"),
        ("5.", "↻", "Vuelvo días después a repasar y profundizar"),
    ]
    top = Inches(2.1)
    for i, (num, ic, txt) in enumerate(steps):
        y = int(top + Inches(0.62) * i)
        s.text(Inches(2.8), y, Inches(8.0), Inches(0.55),
               [para([run(num + "  ", 20, ACCENT, bold=True, font=MONO),
                      run(ic + "   ", 18, ACCENT),
                      run(txt, 20, "DDDDDD")], LEFT)],
               anim="flyL", delay=0.5 + i * 0.15, dur=700)
    s.text(Inches(2.5), Inches(6.55), Inches(8.3), Inches(0.7),
           [para([run("Esto no es estudiar con IA. Es aplicar ciencia del aprendizaje, "
                      "con IA acelerándolo.", 16, AMBER)], CENTER, line=1.2)],
           anim="fade", delay=1.4, dur=1000)
    s.finalize()


# ---- Slide 13 : la ciencia confirma ----------------------------------------
def slide13():
    s = SlideBuilder()
    s.glow(AMBER, 500, 0.50, 0.30, 8)
    s.text(Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.8),
           [para([run("🧪  La ciencia confirma el método", 34, WHITE, bold=True, font=DISPLAY)],
                 CENTER)], anim="fade", delay=0.2)
    s.text(Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.8),
           [para([run("170.000", 48, ACCENT, bold=True, font=DISPLAY),
                  run("  personas", 20, GREY)], CENTER)], anim="fade", delay=0.4, dur=800)
    s.text(Inches(0.8), Inches(2.35), Inches(11.7), Inches(0.5),
           [para([run("Meta-análisis (Hattie & Donoghue, 2021) sobre 242 estudios de "
                      "técnicas de aprendizaje", 15, GREY2)], CENTER)],
           anim="fade", delay=0.5)
    s.text(Inches(0.8), Inches(2.95), Inches(11.7), Inches(0.5),
           [para([run("Las 2 técnicas más efectivas, por encima de todas las demás:",
                      20, SUBTLE)], CENTER)], anim="fade", delay=0.6)
    # No funciona
    s.text(Inches(1.5), Inches(3.9), Inches(5.0), Inches(0.5),
           [para([run("✕  NO FUNCIONA", 17, PRIMARY, bold=True, spc=1.5)], LEFT)],
           anim="fade", delay=0.8)
    nofunc = ["Releer el PDF 10 veces", "Subrayar y resaltar", "Estudio masivo (cramming)"]
    for i, t in enumerate(nofunc):
        s.text(Inches(1.7), int(Inches(4.55) + Inches(0.55) * i), Inches(5.0), Inches(0.5),
               [para([run("•  " + t, 17, GREY)], LEFT)], anim="fade", delay=0.8)
    # Sí funciona
    s.text(Inches(7.0), Inches(3.9), Inches(5.3), Inches(0.5),
           [para([run("✓  SÍ FUNCIONA", 17, AMBER, bold=True, spc=1.5)], LEFT)],
           anim="fade", delay=0.9)
    sifunc = ["Retrieval Practice (práctica de recuperación)",
              "Spaced Repetition (repaso espaciado)",
              "Volver al tema días después"]
    for i, t in enumerate(sifunc):
        s.text(Inches(7.2), int(Inches(4.55) + Inches(0.55) * i), Inches(5.3), Inches(0.5),
               [para([run("•  " + t, 17, "DDDDDD")], LEFT)],
               anim="fade", delay=1.0 + i * 0.15)
    s.finalize()


# ---- Slide 14 : no tengan miedo --------------------------------------------
def slide14():
    s = SlideBuilder()
    s.glow(ACCENT, 600, 0.50, 0.50, 9)
    s.text(Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.6),
           [para([run("No tengan miedo a ", 54, WHITE, bold=True, font=DISPLAY),
                  run("preguntar", 54, ACCENT, bold=True, font=DISPLAY),
                  run(".", 54, WHITE, bold=True, font=DISPLAY)], CENTER, line=1.08)],
           anim="fade", delay=0.3, dur=1100)
    s.bar(int(SLIDE_W / 2 - Inches(1.0)), Inches(4.1), Inches(2.0), Pt(2), ACCENT)
    s.text(Inches(2.0), Inches(4.6), Inches(9.3), Inches(1.2),
           [para([run("Preocúpense solo de entender ", 26, SUBTLE),
                  run("todo a la perfección", 26, WHITE, bold=True),
                  run(".", 26, SUBTLE)], CENTER, line=1.25)],
           anim="fade", delay=1.1)
    s.finalize()


# ---- Slide 15 : ustedes se vuelven más inteligentes ------------------------
def slide15():
    s = SlideBuilder()
    s.glow(AMBER, 500, 0.40, 0.40, 8)
    s.glow(ACCENT, 400, 0.65, 0.65, 8)
    s.text(Inches(1.5), Inches(0.9), Inches(10.3), Inches(0.7),
           [para([run("La IA no los va a hacer más inteligentes sola.", 24, GREY)], CENTER)],
           anim="fade", delay=0.2)
    s.text(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.0),
           [para([run("Ustedes se vuelven más inteligentes", 42, WHITE, bold=True, font=DISPLAY)],
                 CENTER)], anim="wipeU", delay=0.5, dur=1000)
    lines = ["preguntando bien,", "anotando bien,", "y volviendo a profundizar."]
    for i, t in enumerate(lines):
        s.text(Inches(1.5), int(Inches(3.3) + Inches(0.95) * i), Inches(10.3), Inches(0.8),
               [para([run(t, 32, ACCENT, bold=True)], CENTER)],
               anim="flyL", delay=0.8 + i * 0.25, dur=800)
    s.finalize()


# ---- Slide 16 : amplificador (quote) ---------------------------------------
def slide16():
    s = SlideBuilder()
    s.glow(ACCENT, 700, 0.50, 0.50, 10)
    s.text(Inches(0), Inches(1.0), SLIDE_W, Inches(0.8),
           [para([run("❝", 44, ACCENT)], CENTER)], anim="fade", delay=0.3)
    s.text(Inches(1.2), Inches(1.9), Inches(10.9), Inches(1.6),
           [para([run("“Esta tecnología es un ", 44, WHITE, bold=True, font=DISPLAY),
                  run("amplificador", 44, ACCENT, bold=True, font=DISPLAY),
                  run(".”", 44, WHITE, bold=True, font=DISPLAY)], CENTER, line=1.12)],
           anim="fade", delay=0.5, dur=1000)
    s.text(Inches(2.0), Inches(4.4), Inches(9.3), Inches(0.8),
           [para([run("Al que piensa, lo hace pensar más.", 28, "FFFFFF")], CENTER)],
           anim="fade", delay=0.9)
    s.text(Inches(2.0), Inches(5.5), Inches(9.3), Inches(0.9),
           [para([run("Al que delega todo sin chequear, lo hace cada vez más dependiente.",
                      22, "E0E0E8")], CENTER, line=1.2)],
           anim="fade", delay=1.1)
    s.finalize()


# ---- Slide 17 : mejores profesionales --------------------------------------
def slide17():
    s = SlideBuilder()
    s.glow(ACCENT, 500, 0.50, 0.40, 9)
    s.glow(AMBER, 400, 0.30, 0.70, 8)
    s.text(Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.7),
           [para([run("Todos buscamos lo mismo.", 24, GREY, spc=0.5)], CENTER)],
           anim="fade", delay=0.2)
    s.text(Inches(1.2), Inches(2.2), Inches(10.9), Inches(2.4),
           [para([run("Que salgan de acá", 42, WHITE, bold=True, font=DISPLAY)], CENTER, line=1.15),
            para([run("siendo ", 42, WHITE, bold=True, font=DISPLAY),
                  run("mejores profesionales", 42, ACCENT, bold=True, font=DISPLAY)], CENTER, line=1.15),
            para([run("de los que entraron.", 42, WHITE, bold=True, font=DISPLAY)], CENTER, line=1.15)],
           anim="wipeU", delay=0.5, dur=1100)
    s.text(Inches(0), Inches(5.9), SLIDE_W, Inches(0.5),
           [para([run("● ● ● ● ●", 16, ACCENT)], CENTER)], anim="fade", delay=1.0)
    s.finalize()


# ---- Slide 18 : gracias + QR -----------------------------------------------
def slide18():
    s = SlideBuilder()
    s.glow(ACCENT, 600, 0.50, 0.50, 10)
    s.glow(AMBER, 400, 0.20, 0.30, 8)
    s.text(Inches(1.0), Inches(0.7), Inches(11.3), Inches(1.2),
           [para([run("Muchas ", 60, WHITE, bold=True, font=DISPLAY),
                  run("gracias", 60, ACCENT, bold=True, font=DISPLAY)], CENTER)],
           anim="fade", delay=0.3, dur=800)
    s.text(Inches(1.0), Inches(1.95), Inches(11.3), Inches(1.2),
           [para([run("Mauricio Medina", 28, WHITE, bold=True)], CENTER, space_after=4),
            para([run("Desarrollador · CuyoConnect", 17, GREY)], CENTER, space_after=4),
            para([run("📍 Mendoza, Argentina", 15, GREY)], CENTER)],
           anim="fade", delay=0.6)
    s.image(QR, int(SLIDE_W / 2 - Inches(1.25)), Inches(3.5),
            w=Inches(2.5), h=Inches(2.5), anim="fade", delay=1.0, dur=800,
            round_corners=True)
    s.text(Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.7),
           [para([run("@  ", 26, ACCENT, bold=True),
                  run("mauri.h.m", 26, WHITE, bold=True, spc=0.5)], CENTER)],
           anim="fade", delay=1.4)
    s.finalize()


# ===========================================================================
for fn in [slide00, slide01, slide02, slide03, slide04, slide05, slide06,
           slide07, slide08, slide09, slide10, slide11, slide12, slide13,
           slide14, slide15, slide16, slide17, slide18]:
    fn()

prs.save(str(OUT))
print("OK ->", OUT)
print("Slides:", len(prs.slides._sldIdLst))
